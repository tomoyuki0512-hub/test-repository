#!/usr/bin/env python3
"""光ファイバー2ガイドをまとめたパワポを生成する。"""
import math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

JP_FONT = "Meiryo"
NAVY = RGBColor(0x1F, 0x3A, 0x5F)
BLUE = RGBColor(0x2E, 0x6D, 0xB4)
LIGHT = RGBColor(0xED, 0xF3, 0xFA)
GRAY = RGBColor(0x44, 0x44, 0x44)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def set_font(run, size=18, bold=False, color=GRAY, font=JP_FONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def title_slide(title, subtitle, footer):
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide, NAVY)
    # band
    band = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.7), Inches(2.0))
    tf = band.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    set_font(r, 44, True, WHITE)
    p2 = tf.add_paragraph()
    r = p2.add_run(); r.text = subtitle
    set_font(r, 22, False, RGBColor(0xCF, 0xE0, 0xF2))
    f = slide.shapes.add_textbox(Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.6))
    p = f.text_frame.paragraphs[0]
    r = p.add_run(); r.text = footer
    set_font(r, 14, False, RGBColor(0x9F, 0xB6, 0xD0))
    return slide


def section_slide(num, title):
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide, BLUE)
    box = slide.shapes.add_textbox(Inches(0.9), Inches(2.9), Inches(11.5), Inches(1.8))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = num
    set_font(r, 24, True, RGBColor(0xCF, 0xE0, 0xF2))
    p2 = tf.add_paragraph()
    r = p2.add_run(); r.text = title
    set_font(r, 40, True, WHITE)
    return slide


def content_slide(title):
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide, WHITE)
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), SW, Inches(1.05))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tb = bar.text_frame; tb.word_wrap = True
    tb.margin_left = Inches(0.5); tb.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tb.paragraphs[0]
    r = p.add_run(); r.text = title
    set_font(r, 26, True, WHITE)
    return slide


def add_bullets(slide, items, top=1.4, left=0.7, width=12.0, size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.6))
    tf = box.text_frame; tf.word_wrap = True
    first = True
    for text, level in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        bullet = ("• " if level == 0 else "  – ") + text
        r = p.add_run(); r.text = bullet
        set_font(r, size - level * 2, level == 0 and False, GRAY)
        p.space_after = Pt(6)
    return box


def add_table(slide, headers, rows, top=1.45, left=0.6, width=12.1, height=4.8,
              col_widths=None, fsize=14):
    nrows = len(rows) + 1
    ncols = len(headers)
    gtbl = slide.shapes.add_table(nrows, ncols, Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tbl = gtbl.table
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(w)
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        tf = cell.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = h
        set_font(r, fsize, True, WHITE)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT if i % 2 else WHITE
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = val
            set_font(r, fsize, j == 0, NAVY if j == 0 else GRAY)
    return tbl


def _label(slide, text, left, top, width, height, size, color,
           align=PP_ALIGN.LEFT, bold=True):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    set_font(r, size, bold, color)
    return tb


def _no_shadow(shape):
    shape.shadow.inherit = False


def _arrow_tail(connector):
    ln = connector.line._get_or_add_ln()
    ln.append(ln.makeelement(qn('a:tailEnd'),
                             {'type': 'triangle', 'w': 'med', 'len': 'med'}))


def draw_tir_diagram(slide, left, top, width, height):
    """全反射の模式図：横から見たファイバー内を光がジグザグに反射して進む。"""
    clad = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(left), Inches(top), Inches(width), Inches(height))
    clad.fill.solid(); clad.fill.fore_color.rgb = RGBColor(0xCF, 0xE0, 0xF2)
    clad.line.color.rgb = NAVY; clad.line.width = Pt(1.25); _no_shadow(clad)

    core_h = height * 0.46
    core_top = top + (height - core_h) / 2
    core = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(left), Inches(core_top), Inches(width), Inches(core_h))
    core.fill.solid(); core.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xC4)
    core.line.color.rgb = RGBColor(0xD8, 0xB8, 0x40); core.line.width = Pt(0.75)
    _no_shadow(core)

    y_hi, y_lo = core_top, core_top + core_h
    y_mid = core_top + core_h / 2
    n = 4
    xs = [left + width * i / n for i in range(n + 1)]
    pts = [(left, y_mid)]
    for i in range(1, n):
        pts.append((xs[i], y_hi if i % 2 == 1 else y_lo))
    pts.append((left + width, y_mid))
    for i in range(len(pts) - 1):
        x1, y1 = pts[i]; x2, y2 = pts[i + 1]
        c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                       Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        c.line.color.rgb = RGBColor(0xD1, 0x2E, 0x2E); c.line.width = Pt(2.5)
        if i == len(pts) - 2:
            _arrow_tail(c)
    for i in range(1, n):
        mx = xs[i]; my = y_hi if i % 2 == 1 else y_lo
        d = 0.09
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     Inches(mx - d / 2), Inches(my - d / 2), Inches(d), Inches(d))
        dot.fill.solid(); dot.fill.fore_color.rgb = RGBColor(0xD1, 0x2E, 0x2E)
        dot.line.fill.background(); _no_shadow(dot)
    _label(slide, "クラッド（跳ね返す壁）", left + 0.12, top + 0.03, width - 0.24, 0.3,
           11, NAVY, PP_ALIGN.LEFT)
    _label(slide, "コア（光の通り道）", left + 0.12, y_lo + 0.02, width - 0.24, 0.3,
           11, RGBColor(0xB8, 0x90, 0x10), PP_ALIGN.CENTER)


def draw_cross_section(slide, cx, cy, outer_d):
    """断面図：被覆・クラッド・コアの同心円（cx, cy, outer_d はインチ）。"""
    layers = [
        (outer_d, RGBColor(0xDD, 0xE3, 0xEA), NAVY),
        (outer_d * 0.72, RGBColor(0xBF, 0xD6, 0xEE), BLUE),
        (outer_d * 0.34, RGBColor(0xFF, 0xEC, 0x99), RGBColor(0xD8, 0xB8, 0x40)),
    ]
    for d, fill, line in layers:
        o = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                   Inches(cx - d / 2), Inches(cy - d / 2), Inches(d), Inches(d))
        o.fill.solid(); o.fill.fore_color.rgb = fill
        o.line.color.rgb = line; o.line.width = Pt(1.25); _no_shadow(o)


def _legend(slide, left, top, items, sw=0.24, gap=0.46, size=12):
    for i, (color, text) in enumerate(items):
        y = top + i * gap
        sq = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(left), Inches(y), Inches(sw), Inches(sw))
        sq.fill.solid(); sq.fill.fore_color.rgb = color
        sq.line.color.rgb = NAVY; sq.line.width = Pt(0.5); _no_shadow(sq)
        _label(slide, text, left + sw + 0.12, y - 0.02, 4.2, 0.3, size, GRAY,
               PP_ALIGN.LEFT, bold=False)


# ---------- 1. Title ----------
title_slide(
    "光ファイバー 完全ガイド",
    "第1部：基礎のしくみ ／ 第2部：ケーブル・コードと接続部材",
    "出典：docs/optical-fiber-guide.md・docs/optical-fiber-cable-types.md をまとめたスライド版",
)

# ---------- 2. Agenda ----------
s = content_slide("目次（アジェンダ）")
add_bullets(s, [
    ("第1部　光ファイバーの基礎", 0),
    ("光ファイバーとは／なぜ光を使うのか", 1),
    ("全反射のしくみ／3層構造", 1),
    ("シングルモードとマルチモード", 1),
    ("データの流れ・メリデメ・身近な用途", 1),
    ("第2部　ケーブル・コードと接続部材", 0),
    ("テープスロット型 vs 層より型ケーブル", 1),
    ("単心／メガネ／FOコードの違い", 1),
    ("ケーブルとコードの違い", 1),
    ("成端箱とクロージャ／コネクタと規格", 1),
    ("付録：用語集", 0),
], top=1.35, size=18)

# ---------- Section 1 ----------
section_slide("第1部", "光ファイバーの基礎")

# 光ファイバーとは
s = content_slide("光ファイバーとは？（ざっくり結論）")
add_bullets(s, [
    ("髪の毛ほどの細さ（約0.125mm）の「ガラスの糸」", 0),
    ("中を光（点いた/消えた）が通り、1と0でデータを運ぶ", 0),
    ("速くて遠くまで弱まりにくい → 大量・高速・長距離", 0),
    ("家の光回線・海底ケーブル・データセンターの土台", 0),
    ("イメージ：光のモールス信号をガラスのストロー内で高速点滅", 0),
], top=1.6, size=20)

# なぜ光
s = content_slide("なぜ電気でなく「光」を使う？")
add_table(s,
    ["比較", "銅線（電気）", "光ファイバー（光）"],
    [
        ["速さ・容量", "限界がある", "けた違いに大きい"],
        ["遠くへ送る", "すぐ弱まる", "弱まりにくい"],
        ["ノイズ", "電磁波に弱い", "電気ノイズに強い"],
        ["盗聴", "されやすい", "されにくい"],
        ["重さ・太さ", "重く太い", "軽く細い"],
    ],
    col_widths=[3.0, 4.55, 4.55], fsize=16)

# 全反射
s = content_slide("いちばん大事なしくみ：全反射")
add_bullets(s, [
    ("光は密度の違う物質の境目で曲がる/跳ね返る", 0),
    ("ある角度より浅いと100%跳ね返る ＝ 全反射", 0),
    ("たとえ：プールの底から水面を見ると底が鏡のように映る", 1),
    ("中心(コア)を外側(クラッド)で包む構造", 0),
    ("ジグザグに反射しながら進み曲げても漏れにくい", 0),
], top=1.55, left=0.6, width=6.1, size=18)
draw_tir_diagram(s, left=7.0, top=2.3, width=5.7, height=2.3)
_label(s, "▲ 光が全反射しながら進むようす", 7.0, 4.9, 5.7, 0.4, 13, NAVY, PP_ALIGN.CENTER)

# 構造
s = content_slide("3層構造（輪切りにすると）")
add_table(s,
    ["部分", "役割", "たとえ"],
    [
        ["コア（芯）", "光が通る道（屈折率 高）", "水の通り道"],
        ["クラッド", "跳ね返す（屈折率 低）", "鏡の壁"],
        ["被覆", "傷・水・曲げから守る", "服・よろい"],
    ],
    top=1.55, left=0.5, width=7.2, col_widths=[1.7, 3.3, 2.2], fsize=14, height=3.0)
add_bullets(s, [
    ("ポイント：コア＞クラッドの屈折率差が全反射を生む", 0),
], top=4.95, left=0.5, width=7.2, size=16)
_label(s, "断面図（輪切り）", 8.0, 1.5, 4.6, 0.4, 15, NAVY, PP_ALIGN.CENTER)
draw_cross_section(s, cx=10.3, cy=3.35, outer_d=2.6)
_legend(s, 8.55, 5.0, [
    (RGBColor(0xFF, 0xEC, 0x99), "コア（光が通る芯）"),
    (RGBColor(0xBF, 0xD6, 0xEE), "クラッド（跳ね返す層）"),
    (RGBColor(0xDD, 0xE3, 0xEA), "被覆（保護の服）"),
])

# 種類
s = content_slide("種類：シングルモード vs マルチモード")
add_table(s,
    ["", "シングルモード(SMF)", "マルチモード(MMF)"],
    [
        ["コア径", "細い（約9µm）", "太い（約50〜62.5µm）"],
        ["距離", "長距離が得意", "短距離向け"],
        ["容量", "非常に大きい", "大きい"],
        ["コスト", "機器が高め", "比較的安い"],
        ["主用途", "幹線・光回線・海底", "建物内・データセンター内"],
    ],
    col_widths=[2.6, 4.75, 4.75], fsize=16)

# データの流れ
s = content_slide("データの流れ：電気 → 光 → 電気")
add_bullets(s, [
    ("送信側：電気の1/0をレーザー/LEDの点滅に変換", 0),
    ("伝送中：光が全反射しながら高速で進む", 0),
    ("受信側：フォトダイオードが光を電気の1/0へ戻す", 0),
    ("長距離では中継器・光増幅器(EDFA)で光を増幅して届ける", 0),
], top=1.6, size=20)

# メリデメ
s = content_slide("メリットとデメリット")
add_table(s,
    ["メリット", "デメリット"],
    [
        ["超高速・大容量", "折れ・急な曲げに弱い"],
        ["遠くまで届く（低減衰）", "接続作業がデリケート（専用工具）"],
        ["電気ノイズに強い", "電気（電力）は送れない"],
        ["軽くて細い", "端末で電気⇄光の変換機器が必要"],
        ["盗聴されにくい", ""],
    ],
    col_widths=[6.05, 6.05], fsize=16)

# 用途
s = content_slide("身近な使われ方")
add_bullets(s, [
    ("家庭の光回線（FTTH）：電柱→家へ引込み、ONUで光⇄電気変換", 0),
    ("海底ケーブル：大陸間を結び国際通信を支える", 0),
    ("データセンター：サーバー同士を光で高速接続", 0),
    ("携帯基地局間：スマホの裏側も基地局間は光", 0),
    ("医療・工業：内視鏡、レーザー加工、センサーなど", 0),
], top=1.6, size=20)

# ---------- Section 2 ----------
section_slide("第2部", "ケーブル・コードと接続部材")

# テープスロット vs 層より
s = content_slide("テープスロット型 vs 層より型ケーブル")
add_table(s,
    ["観点", "テープスロット型", "層より型（ストランド型）"],
    [
        ["心線の並べ方", "テープ心線を溝(スロット)に積層", "心線をテンションメンバ周りに層状撚り"],
        ["心数・密度", "超多心・高密度が得意", "中小心数向け"],
        ["接続作業", "多心一括融着で速い", "単心ごとに融着"],
        ["単心の分岐", "テープ単位(SZ撚りで中間分岐可)", "単心の識別・分岐がしやすい"],
        ["主用途", "幹線・大容量・地下管路", "中小規模配線・引き落とし"],
    ],
    col_widths=[2.3, 4.9, 4.9], fsize=14)

# コード3種
s = content_slide("単心 / メガネ / FOコードの違い")
add_table(s,
    ["", "心線数", "形・特徴", "主な用途"],
    [
        ["単心コード", "1心", "円1つ。単方向(シンプレクス)", "機器接続・パッチ片側"],
        ["メガネコード", "2心", "8の字型。双方向・裂いて分離", "スイッチ／サーバー間"],
        ["FOコード", "4/8/12心", "テープ心線を1心ずつに展開(ファンアウト)", "架内配線・テープ心線の成端"],
    ],
    col_widths=[2.2, 1.6, 4.8, 3.5], fsize=14, height=3.0)
add_bullets(s, [
    ("FO＝Fiber Optic ではなく「ファンアウト(Fan-Out)」の略", 0),
], top=4.9, size=18)

# ケーブルとコード
s = content_slide("ケーブルとコードの違い")
add_table(s,
    ["観点", "コード（cord）", "ケーブル（cable）"],
    [
        ["主な場所", "屋内・機器のそば", "屋外・長距離・幹線"],
        ["心数", "1〜十数心が中心", "数十〜数千心も"],
        ["強度部材", "アラミド繊維で軽く補強", "鋼線/FRPで頑丈に"],
        ["外被・耐性", "柔らかく取り回し重視", "厚く頑丈・防水耐候"],
        ["役割", "機器⇄配線盤の最後の区間", "区間と区間を結ぶ幹"],
    ],
    col_widths=[2.4, 4.85, 4.85], fsize=15)

# 成端箱とクロージャ
s = content_slide("成端箱とクロージャの違い")
add_table(s,
    ["観点", "成端箱", "クロージャ"],
    [
        ["位置", "ケーブルの端末", "ケーブルの途中"],
        ["主目的", "心線をコネクタ化(成端)して機器へ", "接続・分岐点を保護・収納"],
        ["設置場所", "屋内が中心", "屋外(架空・地中)が中心"],
        ["防水", "屋内前提(屋外型もあり)", "防水・防塵が前提"],
        ["イメージ", "配線盤/パッチパネル", "電線途中の保護ボックス"],
    ],
    col_widths=[2.3, 4.9, 4.9], fsize=15)

# コネクタと規格
s = content_slide("コネクタと規格（つなぎ口の種類）")
add_table(s,
    ["コネクタ", "形・固定方法", "よく見る場所"],
    [
        ["SC", "四角形・押して固定（プッシュプル）", "家庭のONU・通信機器"],
        ["LC", "SCを小型化・ツメで固定", "データセンター・スイッチ"],
        ["ST", "丸型・ひねって固定（バヨネット）", "古めの設備・LAN"],
        ["FC", "丸型・ねじで締めて固定", "計測器・産業用"],
    ],
    top=1.9, left=0.9, width=11.5, col_widths=[1.9, 5.7, 3.9], fsize=17, height=3.8)

# まとめ
s = content_slide("まとめ")
add_bullets(s, [
    ("光ファイバー＝全反射で光を閉じ込めて運ぶガラスの糸", 0),
    ("コア径でシングルモード(長距離)／マルチモード(短距離)", 0),
    ("集合方式：テープスロット型(高密度)／層より型(分岐容易)", 0),
    ("コード：単心(1)・メガネ(2)・FO(テープ展開)で使い分け", 0),
    ("ケーブル=屋外頑丈多心／コード=屋内柔軟機器寄り", 0),
    ("成端箱=端末でコネクタ化／クロージャ=途中を防水保護", 0),
], top=1.5, size=20)

# ---------- 付録：用語集 ----------
section_slide("付録", "用語集")

s = content_slide("用語集 ①（基礎のことば）")
add_table(s,
    ["用語", "意味"],
    [
        ["コア", "光が通る中心の芯（屈折率 高）"],
        ["クラッド", "光をはね返す外側の層（屈折率 低）"],
        ["被覆／ジャケット", "傷・水・曲げから守る外装"],
        ["全反射", "光が境界で100%はね返る現象。伝送の基本原理"],
        ["屈折率", "光の曲がりやすさ。コア＞クラッド"],
        ["シングルモード(SMF)", "細い芯。長距離・大容量向き"],
        ["マルチモード(MMF)", "太い芯。短距離・低コスト向き"],
        ["減衰", "信号が進むうちに弱まること"],
        ["波長", "光の「色」にあたる性質（nmで表す）"],
        ["WDM", "複数波長を1本に多重して容量を増やす技術"],
    ],
    top=1.35, left=0.6, width=12.1, col_widths=[3.4, 8.7], fsize=14, height=5.6)

s = content_slide("用語集 ②（機器・配線のことば）")
add_table(s,
    ["用語", "意味"],
    [
        ["FTTH", "家庭まで光ファイバを引く方式"],
        ["ONU", "家庭で光⇄電気を変換する装置"],
        ["EDFA", "光を直接増幅する装置（光増幅器の一種）"],
        ["テンションメンバ", "ケーブルの抗張力体（鋼線/FRPなどの強度部材）"],
        ["テープ心線", "複数心を並べてテープ状にした心線"],
        ["ファンアウト(FO)", "テープ心線を1心ずつのコネクタに展開すること"],
        ["成端", "心線をコネクタ化して端末を仕上げること"],
        ["ピグテール", "片端にコネクタが付いた短い心線"],
        ["融着接続", "ファイバ同士を熱で溶かして接続する方法"],
        ["クロージャ", "ケーブル途中の接続・分岐点を保護する箱"],
    ],
    top=1.35, left=0.6, width=12.1, col_widths=[3.4, 8.7], fsize=14, height=5.6)

out = "/home/user/test-repository/presentations/optical-fiber-presentation.pptx"
prs.save(out)
print("saved:", out, "slides:", len(prs.slides._sldIdLst))
